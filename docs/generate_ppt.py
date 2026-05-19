"""
Generate PPT presentasi DocuSim.
Output: docs/DocuSim_Presentation.pptx
Run: python docs/generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# Palette
PRIMARY = RGBColor(0x66, 0x7E, 0xEA)
PRIMARY_DARK = RGBColor(0x76, 0x4B, 0xA2)
DARK = RGBColor(0x1E, 0x29, 0x3B)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
GRAY = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x10, 0xB9, 0x81)
DANGER = RGBColor(0xE1, 0x1D, 0x48)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background() if line is None else None
    return s


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def add_bullets(slide, x, y, w, h, items, size=16, color=DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4),
                 subtitle, size=14, color=GRAY)


def footer(slide, page):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.3),
             "DocuSim — Dokumen Similarity dengan TF-IDF + Cosine", size=10, color=GRAY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
             f"{page}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body, color=PRIMARY):
    add_rect(slide, x, y, w, h, LIGHT)
    add_rect(slide, x, y, Inches(0.12), h, color)  # left bar
    add_text(slide, x + Inches(0.3), y + Inches(0.15), w - Inches(0.4), Inches(0.45),
             title, size=15, bold=True, color=DARK)
    add_text(slide, x + Inches(0.3), y + Inches(0.6), w - Inches(0.4), h - Inches(0.7),
             body, size=12, color=GRAY)


# ============================================================================
# SLIDE 1 — Cover
# ============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PRIMARY)
add_rect(s, 0, Inches(5.5), SW, Inches(2), PRIMARY_DARK)
add_text(s, Inches(0.8), Inches(0.6), Inches(2), Inches(0.4),
         "DocuSim", size=18, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.2), Inches(12), Inches(1),
         "Sistem Deteksi Kemiripan", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.0), Inches(12), Inches(1),
         "Abstrak Dokumen", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(4.2), Inches(12), Inches(0.6),
         "Berbasis Web dengan Laravel + TF-IDF + Cosine Similarity",
         size=20, color=WHITE)
add_text(s, Inches(0.8), Inches(5.9), Inches(12), Inches(0.4),
         "Mata Kuliah Pemrograman Web Lanjut", size=14, color=WHITE)
add_text(s, Inches(0.8), Inches(6.3), Inches(12), Inches(0.4),
         "Semester 6", size=14, color=WHITE)

# ============================================================================
# SLIDE 2 — Latar Belakang
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Latar Belakang")
add_bullets(s, Inches(0.7), Inches(1.3), Inches(12), Inches(5), [
    "Plagiarisme akademik makin sulit dideteksi manual seiring banyaknya "
    "skripsi/tesis/jurnal yang diunggah setiap tahunnya.",
    "Dosen pembimbing perlu cara cepat untuk memeriksa apakah abstrak yang "
    "diajukan mahasiswa mirip dengan dokumen yang sudah ada di repositori.",
    "Algoritma TF-IDF + Cosine Similarity terbukti efektif untuk mengukur "
    "kemiripan dokumen teks pendek seperti abstrak.",
    "Dibutuhkan sistem berbasis web yang ringan, mudah dipakai, dan bisa "
    "menampilkan hasil pengecekan secara visual.",
], size=18)
footer(s, 2)

# ============================================================================
# SLIDE 3 — Rumusan Masalah & Tujuan
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Rumusan Masalah & Tujuan")
add_text(s, Inches(0.7), Inches(1.2), Inches(6), Inches(0.5),
         "Rumusan Masalah", size=18, bold=True, color=PRIMARY_DARK)
add_bullets(s, Inches(0.7), Inches(1.7), Inches(6), Inches(3.5), [
    "Bagaimana menerapkan TF-IDF + Cosine Similarity untuk mendeteksi "
    "kemiripan abstrak?",
    "Bagaimana mengintegrasikan model NLP ke dalam sistem berbasis Laravel?",
    "Bagaimana menampilkan hasil pengecekan secara informatif dan visual?",
], size=15)
add_text(s, Inches(7), Inches(1.2), Inches(6), Inches(0.5),
         "Tujuan", size=18, bold=True, color=PRIMARY_DARK)
add_bullets(s, Inches(7), Inches(1.7), Inches(5.8), Inches(3.5), [
    "Membangun web Laravel yang bisa menerima abstrak sebagai input.",
    "Menerapkan algoritma TF-IDF + Cosine Similarity untuk pembobotan "
    "dan pengukuran kemiripan.",
    "Menyajikan hasil dengan visualisasi (gauge, bar chart, top-terms).",
    "Menyediakan CRUD corpus + import CSV agar data referensi mudah "
    "dikelola.",
], size=15)
footer(s, 3)

# ============================================================================
# SLIDE 4 — Stack Teknologi
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Stack Teknologi")
items = [
    ("Backend", "Laravel 11 (PHP 8.2)\nEloquent ORM, Blade, Symfony Process",
     PRIMARY),
    ("Frontend", "Tailwind CSS v3 (CDN)\nChart.js untuk visualisasi data",
     ACCENT),
    ("Database", "SQLite (default)\nKolom JSON untuk hasil & top-terms",
     PRIMARY_DARK),
    ("Model NLP", "TF-IDF + Cosine Similarity\nDriver PHP native (default)\n"
     "Driver Python scikit-learn (opsional)", DANGER),
]
x = Inches(0.5)
for i, (t, b, c) in enumerate(items):
    card(s, x + Inches(i * 3.15), Inches(1.4), Inches(3), Inches(2.6),
         t, b, color=c)
add_text(s, Inches(0.7), Inches(4.4), Inches(12), Inches(0.5),
         "Mengapa Multi-Driver?", size=18, bold=True, color=PRIMARY_DARK)
add_bullets(s, Inches(0.7), Inches(4.9), Inches(12), Inches(2.0), [
    "Driver PHP native — tanpa instalasi tambahan, jalan langsung di "
    "shared hosting biasa.",
    "Driver Python — siap diganti ke model deep-learning seperti IndoBERT "
    "atau Sentence-Transformers tanpa mengubah kode Laravel.",
    "Switch driver hanya dengan ubah satu baris di .env (SIMILARITY_DRIVER).",
], size=14)
footer(s, 4)

# ============================================================================
# SLIDE 5 — Arsitektur Sistem
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Arsitektur Sistem")
# Boxes
def archbox(x, y, w, h, title, sub, color):
    add_rect(s, x, y, w, h, color)
    add_text(s, x, y + Inches(0.15), w, Inches(0.4), title,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.65), w, Inches(0.6), sub,
             size=11, color=WHITE, align=PP_ALIGN.CENTER)

archbox(Inches(0.6), Inches(2.5), Inches(2.6), Inches(1.5),
        "Browser", "Blade view\nTailwind + Chart.js", PRIMARY)
archbox(Inches(4.1), Inches(2.5), Inches(2.6), Inches(1.5),
        "Laravel", "Controller + Service\n+ Eloquent ORM", PRIMARY_DARK)
archbox(Inches(7.6), Inches(2.5), Inches(2.6), Inches(1.5),
        "Engine TF-IDF", "PHP native (default)\nPython opsional", DANGER)
archbox(Inches(11.0), Inches(2.5), Inches(1.8), Inches(1.5),
        "SQLite", "Corpus + History", ACCENT)

# Arrows (simple lines)
def arrow(x1, x2, y, label):
    s.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
    add_text(s, Inches(x1), Inches(y - 0.4), Inches(x2 - x1), Inches(0.3),
             label, size=10, color=GRAY, align=PP_ALIGN.CENTER)

arrow(3.2, 4.1, 3.25, "HTTP")
arrow(6.7, 7.6, 3.25, "panggil")
arrow(8.0, 11.0, 3.25, "Eloquent")

add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.5),
         "Alur Pengecekan Similarity", size=18, bold=True, color=PRIMARY_DARK)
add_bullets(s, Inches(0.7), Inches(5.0), Inches(12), Inches(2.0), [
    "User submit abstrak → Controller validasi → Service ambil corpus "
    "dari DB.",
    "Engine TF-IDF tokenize, hitung DF & IDF, bentuk vector + l2-normalize, "
    "lalu cosine similarity.",
    "Hasil ranking + top-terms disimpan ke tabel similarity_checks dan "
    "ditampilkan dengan Chart.js.",
], size=14)
footer(s, 5)

# ============================================================================
# SLIDE 6 — Algoritma TF-IDF
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Algoritma: TF-IDF + Cosine Similarity")
add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.5),
         "1. Term Frequency (TF)", size=16, bold=True, color=PRIMARY_DARK)
add_text(s, Inches(0.9), Inches(1.65), Inches(12), Inches(0.5),
         "tf(t, d) = frekuensi kata t di dokumen d", size=14, color=DARK)

add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.5),
         "2. Inverse Document Frequency (IDF) — Smooth", size=16, bold=True,
         color=PRIMARY_DARK)
add_text(s, Inches(0.9), Inches(2.75), Inches(12), Inches(0.5),
         "idf(t) = ln( (1 + N) / (1 + df(t)) ) + 1", size=14, color=DARK)
add_text(s, Inches(0.9), Inches(3.15), Inches(12), Inches(0.4),
         "N = total dokumen,  df(t) = jumlah dokumen yang memuat kata t",
         size=12, color=GRAY)

add_text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.5),
         "3. Bobot TF-IDF + Normalisasi L2", size=16, bold=True,
         color=PRIMARY_DARK)
add_text(s, Inches(0.9), Inches(4.15), Inches(12), Inches(0.5),
         "tfidf(t, d) = tf(t, d) × idf(t)", size=14, color=DARK)
add_text(s, Inches(0.9), Inches(4.55), Inches(12), Inches(0.5),
         "Vector dinormalisasi: v / ||v||₂", size=14, color=DARK)

add_text(s, Inches(0.7), Inches(5.1), Inches(12), Inches(0.5),
         "4. Cosine Similarity", size=16, bold=True, color=PRIMARY_DARK)
add_text(s, Inches(0.9), Inches(5.55), Inches(12), Inches(0.5),
         "cos(A, B) = (A · B) / (||A|| × ||B||)", size=14, color=DARK)
add_text(s, Inches(0.9), Inches(5.95), Inches(12), Inches(0.4),
         "Karena vector sudah dinormalisasi → cos = dot(A, B). "
         "Skor 0..1, makin tinggi makin mirip.",
         size=12, color=GRAY)
footer(s, 6)

# ============================================================================
# SLIDE 7 — Pre-processing
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Pre-processing Teks Bahasa Indonesia")
steps = [
    ("Case Folding", "Ubah seluruh teks menjadi huruf kecil agar 'Sistem' "
     "dan 'sistem' dianggap sama."),
    ("Cleaning", "Hapus tanda baca, angka, dan karakter non-alfabet."),
    ("Tokenisasi", "Pecah teks menjadi unit kata."),
    ("Stopword Removal", "Hilangkan kata umum Bahasa Indonesia (yang, "
     "untuk, pada, dengan, dll) dengan kamus Sastrawi."),
    ("N-gram (1-2)", "Bentuk unigram + bigram untuk menangkap frasa "
     "(mis. 'tf idf', 'cosine similarity')."),
]
for i, (t, b) in enumerate(steps):
    y = Inches(1.3 + i * 1.05)
    add_rect(s, Inches(0.7), y, Inches(0.6), Inches(0.6), PRIMARY)
    add_text(s, Inches(0.7), y + Inches(0.1), Inches(0.6), Inches(0.4),
             str(i + 1), size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), y, Inches(11.5), Inches(0.4),
             t, size=15, bold=True, color=DARK)
    add_text(s, Inches(1.5), y + Inches(0.4), Inches(11.5), Inches(0.55),
             b, size=12, color=GRAY)
footer(s, 7)

# ============================================================================
# SLIDE 8 — Fitur Sistem
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Fitur Sistem")
features = [
    ("🎯  Cek Similarity",
     "Input abstrak, sistem cari dokumen paling mirip dengan ranking top-K."),
    ("📊  Visualisasi",
     "Gauge skor tertinggi, bar chart ranking, top-terms TF-IDF."),
    ("🗂️  CRUD Corpus",
     "Kelola dokumen referensi: judul, penulis, tahun, kategori, abstrak."),
    ("📥  Import CSV",
     "Tambah corpus batch dengan format minimal no, judul, abstrak."),
    ("🕘  Riwayat",
     "Setiap pengecekan tersimpan beserta hasilnya di database."),
    ("🎨  UI Modern",
     "Sidebar gradient, kartu interaktif, badge level, fully responsive."),
]
for i, (t, b) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.3 + row * 1.85)
    card(s, x, y, Inches(6), Inches(1.6), t, b, color=PRIMARY)
footer(s, 8)

# ============================================================================
# SLIDE 9 — Database Schema
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Skema Database")

# Tabel 1
add_rect(s, Inches(0.7), Inches(1.3), Inches(5.8), Inches(0.6), PRIMARY)
add_text(s, Inches(0.9), Inches(1.4), Inches(5.5), Inches(0.4),
         "corpus_documents", size=16, bold=True, color=WHITE)
fields1 = ["id (PK)", "title", "author", "year", "category",
           "abstract (longtext)", "created_at, updated_at"]
for i, f in enumerate(fields1):
    add_rect(s, Inches(0.7), Inches(1.9 + i * 0.42),
             Inches(5.8), Inches(0.42), LIGHT)
    add_text(s, Inches(0.95), Inches(1.95 + i * 0.42),
             Inches(5.5), Inches(0.4), f, size=12, color=DARK)

# Tabel 2
add_rect(s, Inches(7.0), Inches(1.3), Inches(5.7), Inches(0.6), PRIMARY_DARK)
add_text(s, Inches(7.2), Inches(1.4), Inches(5.4), Inches(0.4),
         "similarity_checks", size=16, bold=True, color=WHITE)
fields2 = ["id (PK)", "input_title", "input_abstract",
           "highest_score (0..1)", "results (JSON ranking)",
           "top_terms (JSON bobot kata)", "created_at, updated_at"]
for i, f in enumerate(fields2):
    add_rect(s, Inches(7.0), Inches(1.9 + i * 0.42),
             Inches(5.7), Inches(0.42), LIGHT)
    add_text(s, Inches(7.25), Inches(1.95 + i * 0.42),
             Inches(5.4), Inches(0.4), f, size=12, color=DARK)

add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.5),
         "Penyimpanan JSON memudahkan render visualisasi tanpa "
         "menghitung ulang.", size=13, color=GRAY)
footer(s, 9)

# ============================================================================
# SLIDE 10 — Visualisasi
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Visualisasi Hasil")
viz = [
    ("Gauge Chart",
     "Doughnut chart yang menampilkan skor kemiripan tertinggi.\n"
     "Warna mengikuti tingkat: hijau (rendah), kuning (sedang), "
     "merah (tinggi).", ACCENT),
    ("Bar Chart Ranking",
     "Horizontal bar chart yang membandingkan persentase kemiripan "
     "Top-K dokumen.\nMudah dibandingkan secara visual.", PRIMARY),
    ("Top Terms TF-IDF",
     "Bar chart bobot kata terpenting dari abstrak input.\n"
     "Menunjukkan kata mana yang paling berkontribusi pada hasil.",
     PRIMARY_DARK),
]
for i, (t, b, c) in enumerate(viz):
    card(s, Inches(0.6 + i * 4.25), Inches(1.4),
         Inches(4.05), Inches(2.8), t, b, color=c)

add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.5),
         "Mengapa Visualisasi Penting?", size=18, bold=True,
         color=PRIMARY_DARK)
add_bullets(s, Inches(0.7), Inches(5.0), Inches(12), Inches(2.0), [
    "Membantu user memahami hasil tanpa membaca angka mentah.",
    "Warna otomatis (hijau/kuning/merah) memberikan keputusan cepat.",
    "Top-terms menunjukkan kata yang membuat dokumen mirip — bisa "
    "memvalidasi hasil secara intuitif.",
], size=14)
footer(s, 10)

# ============================================================================
# SLIDE 11 — Demo Flow
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Alur Demo Aplikasi")
flow = [
    ("Buka Dashboard",
     "Tampilan statistik: jumlah corpus, jumlah pengecekan, rata-rata skor."),
    ("Tambah Corpus",
     "Input manual atau import CSV dengan format no, judul, abstrak."),
    ("Cek Similarity",
     "Tempel abstrak → klik Analisis → tunggu beberapa detik."),
    ("Lihat Hasil",
     "Gauge + bar chart + top-terms muncul. Detail ranking dengan "
     "expand abstrak."),
    ("Riwayat",
     "Buka kembali pengecekan sebelumnya, atau hapus dari riwayat."),
]
for i, (t, b) in enumerate(flow):
    y = Inches(1.3 + i * 1.1)
    add_rect(s, Inches(0.7), y, Inches(0.7), Inches(0.7), PRIMARY)
    add_text(s, Inches(0.7), y + Inches(0.15), Inches(0.7), Inches(0.4),
             str(i + 1), size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.6), y + Inches(0.05), Inches(11.5), Inches(0.4),
             t, size=15, bold=True, color=DARK)
    add_text(s, Inches(1.6), y + Inches(0.45), Inches(11.5), Inches(0.5),
             b, size=12, color=GRAY)
footer(s, 11)

# ============================================================================
# SLIDE 12 — Kesimpulan
# ============================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Kesimpulan")
add_bullets(s, Inches(0.7), Inches(1.3), Inches(12), Inches(3.5), [
    "Sistem berhasil mengintegrasikan algoritma TF-IDF + Cosine Similarity "
    "ke dalam web Laravel untuk mendeteksi kemiripan abstrak.",
    "Pendekatan multi-driver (PHP native + Python) membuat sistem "
    "fleksibel: ringan saat di-deploy, mudah di-upgrade ke model "
    "deep-learning.",
    "Visualisasi (gauge, bar chart, top-terms) memudahkan user memahami "
    "hasil pengecekan secara intuitif.",
    "Fitur import CSV mempercepat pengisian corpus dalam jumlah besar "
    "tanpa harus input satu per satu.",
], size=16)

add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5),
         "Rencana Pengembangan", size=18, bold=True, color=PRIMARY_DARK)
add_bullets(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.5), [
    "Tambahkan model semantik berbasis IndoBERT untuk menangani parafrase.",
    "Highlight kalimat yang mirip antar dokumen (string matching).",
    "Multi-user dengan autentikasi & role-based access control.",
], size=14)
footer(s, 12)

# ============================================================================
# SLIDE 13 — Terima Kasih
# ============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PRIMARY)
add_rect(s, 0, Inches(5.0), SW, Inches(2.5), PRIMARY_DARK)
add_text(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.2),
         "Terima Kasih", size=72, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.6),
         "Pertanyaan & Diskusi", size=24, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.4),
         "DocuSim — Dokumen Similarity dengan TF-IDF + Cosine Similarity",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
         "Laravel 11 · PHP 8.2 · Tailwind · Chart.js",
         size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================================
# Save
# ============================================================================
out = Path(__file__).parent / "DocuSim_Presentation.pptx"
prs.save(str(out))
print(f"OK -> {out}  ({len(prs.slides)} slides)")
